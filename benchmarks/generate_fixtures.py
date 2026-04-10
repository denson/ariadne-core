"""Generate benchmark fixture documents across multiple formats.

Creates 30+ test documents in benchmarks/fixtures/ for performance testing.
Requires: python-docx, openpyxl, python-pptx (all pulled in by markitdown[all]).

Usage:
    python benchmarks/generate_fixtures.py
"""

from __future__ import annotations

import csv
import json
import random
import textwrap
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"

# Sample content for generating realistic documents
LOREM = (
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do eiusmod "
    "tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim "
    "veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea "
    "commodo consequat. Duis aute irure dolor in reprehenderit in voluptate "
    "velit esse cillum dolore eu fugiat nulla pariatur."
)

TOPICS = [
    ("Quarterly Revenue Report", "finance", "revenue"),
    ("Machine Learning Pipeline Architecture", "engineering", "ml"),
    ("Employee Onboarding Handbook", "hr", "onboarding"),
    ("API Security Audit Findings", "security", "audit"),
    ("Product Roadmap Q4 2025", "product", "roadmap"),
    ("Customer Support Metrics Dashboard", "support", "metrics"),
    ("Database Migration Plan", "engineering", "migration"),
    ("Marketing Campaign Analysis", "marketing", "campaign"),
    ("Incident Response Playbook", "operations", "incident"),
    ("Research Paper: Distributed Systems", "research", "distributed"),
    ("Supply Chain Optimization Report", "logistics", "supply-chain"),
    ("Legal Compliance Review", "legal", "compliance"),
    ("Training Materials: Cloud Infrastructure", "education", "cloud"),
    ("Board Meeting Minutes", "executive", "minutes"),
    ("Technical Specification: Auth Service", "engineering", "auth"),
]

SECTIONS = [
    "Executive Summary",
    "Background",
    "Methodology",
    "Key Findings",
    "Analysis",
    "Recommendations",
    "Implementation Plan",
    "Timeline",
    "Budget Considerations",
    "Risk Assessment",
    "Conclusion",
    "Appendix",
]


def _paragraphs(n: int = 3) -> str:
    """Generate n paragraphs of filler text."""
    paras = []
    for _ in range(n):
        words = LOREM.split()
        random.shuffle(words)
        paras.append(" ".join(words))
    return "\n\n".join(paras)


def _section_content(heading: str, depth: int = 2) -> str:
    """Generate a markdown section with heading and content."""
    prefix = "#" * depth
    return f"{prefix} {heading}\n\n{_paragraphs(random.randint(2, 5))}\n"


def generate_txt_files() -> list[Path]:
    """Generate plain text documents."""
    files = []
    for i, (title, domain, tag) in enumerate(TOPICS[:5]):
        content = f"# {title}\n\n"
        content += f"Domain: {domain} | Tag: {tag}\n\n"
        for section in random.sample(SECTIONS, random.randint(4, 8)):
            content += _section_content(section)
            content += "\n"
        path = FIXTURES_DIR / f"doc_{i+1:02d}_{tag}.txt"
        path.write_text(content, encoding="utf-8")
        files.append(path)
    return files


def generate_html_files() -> list[Path]:
    """Generate HTML documents with various structures."""
    files = []
    for i, (title, domain, tag) in enumerate(TOPICS[5:10]):
        sections_html = ""
        for section in random.sample(SECTIONS, random.randint(3, 6)):
            paragraphs = "\n".join(
                f"    <p>{p}</p>" for p in _paragraphs(random.randint(2, 4)).split("\n\n")
            )
            sections_html += f"  <h2>{section}</h2>\n{paragraphs}\n"

        # Add a table to some
        if i % 2 == 0:
            sections_html += textwrap.dedent("""\
              <h2>Data Summary</h2>
              <table>
                <tr><th>Metric</th><th>Value</th><th>Change</th></tr>
                <tr><td>Revenue</td><td>$1.2M</td><td>+12%</td></tr>
                <tr><td>Users</td><td>45,000</td><td>+8%</td></tr>
                <tr><td>NPS Score</td><td>72</td><td>+5</td></tr>
                <tr><td>Response Time</td><td>1.2s</td><td>-15%</td></tr>
              </table>
            """)

        # Add a list to some
        if i % 2 == 1:
            items = "\n".join(f"    <li>Action item {j+1}: {LOREM[:60]}</li>" for j in range(5))
            sections_html += f"  <h2>Action Items</h2>\n  <ul>\n{items}\n  </ul>\n"

        html = textwrap.dedent(f"""\
            <!DOCTYPE html>
            <html lang="en">
            <head>
              <meta charset="utf-8">
              <title>{title}</title>
            </head>
            <body>
              <h1>{title}</h1>
              <p><em>Domain: {domain} | Category: {tag}</em></p>
            {sections_html}
            </body>
            </html>
        """)
        path = FIXTURES_DIR / f"doc_{i+6:02d}_{tag}.html"
        path.write_text(html, encoding="utf-8")
        files.append(path)
    return files


def generate_csv_files() -> list[Path]:
    """Generate CSV data files."""
    files = []

    # Sales data
    rows = [["Date", "Product", "Region", "Units", "Revenue", "Cost"]]
    products = ["Widget A", "Widget B", "Service X", "Service Y"]
    regions = ["North", "South", "East", "West"]
    for month in range(1, 13):
        for product in products:
            region = random.choice(regions)
            units = random.randint(100, 5000)
            revenue = round(units * random.uniform(10, 100), 2)
            cost = round(revenue * random.uniform(0.3, 0.7), 2)
            rows.append([f"2025-{month:02d}-01", product, region, units, revenue, cost])

    path = FIXTURES_DIR / "doc_11_sales_data.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        csv.writer(f).writerows(rows)
    files.append(path)

    # Server metrics
    rows = [["Timestamp", "Server", "CPU%", "Memory%", "Disk%", "Requests/s", "Errors"]]
    servers = ["web-01", "web-02", "api-01", "api-02", "db-01"]
    for hour in range(24):
        for server in servers:
            rows.append([
                f"2025-10-15T{hour:02d}:00:00Z",
                server,
                round(random.uniform(10, 95), 1),
                round(random.uniform(20, 85), 1),
                round(random.uniform(15, 75), 1),
                random.randint(50, 2000),
                random.randint(0, 50),
            ])
    path = FIXTURES_DIR / "doc_12_server_metrics.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        csv.writer(f).writerows(rows)
    files.append(path)

    return files


def generate_json_files() -> list[Path]:
    """Generate JSON data files."""
    files = []

    # API response mock
    data = {
        "api_version": "2.1",
        "endpoints": [
            {
                "path": f"/api/v2/{name}",
                "method": method,
                "description": f"{name.title()} endpoint for {LOREM[:40]}",
                "auth_required": random.choice([True, False]),
                "rate_limit": random.randint(100, 1000),
                "response_codes": [200, 400, 401, 404, 500],
            }
            for name, method in [
                ("users", "GET"), ("users", "POST"), ("documents", "GET"),
                ("documents", "POST"), ("search", "POST"), ("collections", "GET"),
                ("health", "GET"), ("stats", "GET"),
            ]
        ],
        "schemas": {
            "User": {"id": "uuid", "name": "string", "email": "string"},
            "Document": {"id": "uuid", "title": "string", "content": "text"},
        },
    }
    path = FIXTURES_DIR / "doc_13_api_spec.json"
    path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    files.append(path)

    return files


def generate_markdown_files() -> list[Path]:
    """Generate Markdown documents with rich formatting."""
    files = []

    for i, (title, domain, tag) in enumerate(TOPICS[10:15]):
        content = f"# {title}\n\n"
        content += f"> **Domain:** {domain} | **Category:** {tag}\n\n"
        content += "---\n\n"

        for j, section in enumerate(random.sample(SECTIONS, random.randint(5, 9))):
            content += _section_content(section)

            # Add code blocks to some sections
            if j == 1:
                content += "\n```python\ndef process_document(uri: str) -> dict:\n"
                content += '    """Extract and process a document."""\n'
                content += "    result = extractor.extract(uri)\n"
                content += "    chunks = chunker.chunk(result.markdown)\n"
                content += "    return {\"chunks\": len(chunks)}\n```\n\n"

            # Add tables to some sections
            if j == 3:
                content += "\n| Component | Status | Priority |\n"
                content += "|-----------|--------|----------|\n"
                for k in range(5):
                    status = random.choice(["Complete", "In Progress", "Planned"])
                    priority = random.choice(["High", "Medium", "Low"])
                    content += f"| Component {k+1} | {status} | {priority} |\n"
                content += "\n"

            # Add bullet lists
            if j == 2:
                content += "\nKey points:\n\n"
                for k in range(4):
                    content += f"- Point {k+1}: {LOREM[:80]}\n"
                content += "\n"

        path = FIXTURES_DIR / f"doc_{i+14:02d}_{tag}.md"
        path.write_text(content, encoding="utf-8")
        files.append(path)
    return files


def generate_docx_files() -> list[Path]:
    """Generate DOCX documents using python-docx."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError:
        print("python-docx not available, skipping DOCX generation")
        return []

    files = []
    docx_topics = [
        ("Project Status Report", "management"),
        ("Technical Design Document", "engineering"),
        ("Meeting Notes: Architecture Review", "notes"),
        ("Policy Update: Data Retention", "policy"),
        ("Vendor Evaluation Summary", "procurement"),
    ]

    for i, (title, tag) in enumerate(docx_topics):
        doc = Document()
        doc.add_heading(title, level=0)
        doc.add_paragraph(f"Category: {tag} | Date: 2025-10-15")
        doc.add_paragraph("")

        for section in random.sample(SECTIONS, random.randint(4, 7)):
            doc.add_heading(section, level=1)
            for _ in range(random.randint(2, 4)):
                words = LOREM.split()
                random.shuffle(words)
                doc.add_paragraph(" ".join(words))

            # Add a table to some sections
            if random.random() > 0.6:
                table = doc.add_table(rows=4, cols=3)
                table.style = "Table Grid"
                headers = table.rows[0].cells
                headers[0].text = "Item"
                headers[1].text = "Value"
                headers[2].text = "Status"
                for row_idx in range(1, 4):
                    row = table.rows[row_idx].cells
                    row[0].text = f"Item {row_idx}"
                    row[1].text = str(random.randint(10, 100))
                    row[2].text = random.choice(["Active", "Pending", "Done"])

        path = FIXTURES_DIR / f"doc_{i+19:02d}_{tag}.docx"
        doc.save(str(path))
        files.append(path)
    return files


def generate_pptx_files() -> list[Path]:
    """Generate PPTX presentations using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx not available, skipping PPTX generation")
        return []

    files = []
    pptx_topics = [
        ("Q4 Business Review", "business"),
        ("Engineering All-Hands", "engineering"),
        ("Product Launch Plan", "product"),
    ]

    for i, (title, tag) in enumerate(pptx_topics):
        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Ariadne Core Benchmark | {tag}"

        # Content slides
        for j, section in enumerate(random.sample(SECTIONS, random.randint(6, 10))):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = LOREM[:150]
            for k in range(3):
                p = tf.add_paragraph()
                p.text = f"Point {k+1}: {LOREM[:80]}"
                p.level = 1

        path = FIXTURES_DIR / f"doc_{i+24:02d}_{tag}.pptx"
        prs.save(str(path))
        files.append(path)
    return files


def generate_xlsx_files() -> list[Path]:
    """Generate XLSX spreadsheets using openpyxl."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("openpyxl not available, skipping XLSX generation")
        return []

    files = []

    # Financial data
    wb = Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Month", "Product A", "Product B", "Product C", "Total"])
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
              "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    for month in months:
        a, b, c = (random.randint(10000, 50000) for _ in range(3))
        ws.append([month, a, b, c, a + b + c])

    # Add a second sheet
    ws2 = wb.create_sheet("Expenses")
    ws2.append(["Category", "Q1", "Q2", "Q3", "Q4"])
    categories = ["Salaries", "Infrastructure", "Marketing", "R&D", "Operations"]
    for cat in categories:
        ws2.append([cat] + [random.randint(50000, 200000) for _ in range(4)])

    path = FIXTURES_DIR / "doc_27_financials.xlsx"
    wb.save(str(path))
    files.append(path)

    # Inventory data
    wb = Workbook()
    ws = wb.active
    ws.title = "Inventory"
    ws.append(["SKU", "Name", "Category", "Quantity", "Unit Price", "Warehouse"])
    warehouses = ["US-East", "US-West", "EU-Central"]
    for j in range(50):
        ws.append([
            f"SKU-{j+1:04d}",
            f"Product {j+1}",
            random.choice(["Electronics", "Furniture", "Supplies", "Software"]),
            random.randint(0, 500),
            round(random.uniform(5, 500), 2),
            random.choice(warehouses),
        ])
    path = FIXTURES_DIR / "doc_28_inventory.xlsx"
    wb.save(str(path))
    files.append(path)

    return files


def generate_mixed_content_files() -> list[Path]:
    """Generate files with mixed content types."""
    files = []

    # Large text file (many sections)
    content = "# Comprehensive Technical Reference\n\n"
    for i in range(20):
        content += f"## Chapter {i+1}: Topic {i+1}\n\n"
        content += _paragraphs(random.randint(3, 6))
        content += "\n\n"
        if i % 3 == 0:
            content += "```\n# Code example\nfor i in range(100):\n    print(f'Item {i}')\n```\n\n"
        if i % 4 == 0:
            content += "| Col A | Col B | Col C |\n|-------|-------|-------|\n"
            for j in range(3):
                content += f"| {random.randint(1,100)} | {random.randint(1,100)} | {random.randint(1,100)} |\n"
            content += "\n"
    path = FIXTURES_DIR / "doc_29_large_reference.txt"
    path.write_text(content, encoding="utf-8")
    files.append(path)

    # Minimal document
    path = FIXTURES_DIR / "doc_30_minimal.txt"
    path.write_text("Hello, world.\n", encoding="utf-8")
    files.append(path)

    # Log-like file (no headings)
    lines = []
    for i in range(200):
        level = random.choice(["INFO", "WARN", "ERROR", "DEBUG"])
        msg = random.choice([
            "Request processed successfully",
            "Database connection pool at 80%",
            "Cache miss for key user:12345",
            "Retry attempt 2/3 for external API",
            "Background job completed in 1.2s",
            "Memory usage: 2.1GB / 4.0GB",
        ])
        lines.append(f"2025-10-15T{random.randint(0,23):02d}:{random.randint(0,59):02d}:{random.randint(0,59):02d}Z [{level}] {msg}")
    path = FIXTURES_DIR / "doc_31_application.log"
    path.write_text("\n".join(lines), encoding="utf-8")
    files.append(path)

    # HTML with deeply nested structure
    nested = "<div>\n"
    for i in range(5):
        nested += f'  <div class="section-{i}">\n'
        nested += f"    <h2>Section {i+1}</h2>\n"
        nested += f"    <p>{LOREM}</p>\n"
        for j in range(3):
            nested += f'    <div class="subsection">\n'
            nested += f"      <h3>Subsection {i+1}.{j+1}</h3>\n"
            nested += f"      <p>{LOREM[:100]}</p>\n"
            nested += "    </div>\n"
        nested += "  </div>\n"
    nested += "</div>"

    html = f"""<!DOCTYPE html>
<html>
<head><title>Complex Nested Document</title></head>
<body>
<h1>Complex Nested Document</h1>
{nested}
</body>
</html>"""
    path = FIXTURES_DIR / "doc_32_nested.html"
    path.write_text(html, encoding="utf-8")
    files.append(path)

    return files


def main() -> None:
    """Generate all benchmark fixtures."""
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    random.seed(42)  # Reproducible

    all_files: list[Path] = []
    all_files.extend(generate_txt_files())
    print(f"  Generated {len(all_files)} TXT files")

    html_files = generate_html_files()
    all_files.extend(html_files)
    print(f"  Generated {len(html_files)} HTML files")

    csv_files = generate_csv_files()
    all_files.extend(csv_files)
    print(f"  Generated {len(csv_files)} CSV files")

    json_files = generate_json_files()
    all_files.extend(json_files)
    print(f"  Generated {len(json_files)} JSON files")

    md_files = generate_markdown_files()
    all_files.extend(md_files)
    print(f"  Generated {len(md_files)} Markdown files")

    docx_files = generate_docx_files()
    all_files.extend(docx_files)
    print(f"  Generated {len(docx_files)} DOCX files")

    pptx_files = generate_pptx_files()
    all_files.extend(pptx_files)
    print(f"  Generated {len(pptx_files)} PPTX files")

    xlsx_files = generate_xlsx_files()
    all_files.extend(xlsx_files)
    print(f"  Generated {len(xlsx_files)} XLSX files")

    mixed_files = generate_mixed_content_files()
    all_files.extend(mixed_files)
    print(f"  Generated {len(mixed_files)} mixed-content files")

    print(f"\nTotal: {len(all_files)} benchmark fixtures in {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
